"""Document parsing: a file on disk -> a list of text units, one per "page".

Every format funnels into the same `list[str]` contract that `chunk_pages` already
consumes, so adding a format costs one extractor and nothing downstream changes.

WHAT A "PAGE" MEANS PER FORMAT — this is the honest caveat. Only PDF has real
pages. For the others we emit the closest natural unit the format actually has:

    .pdf          a real page
    .pptx         a slide            (natural, exact)
    .xlsx         a worksheet        (natural, exact)
    .docx         a *section* — split on the document's own explicit page breaks
                  where the author inserted them, otherwise on a size budget
    .txt / .md    a section on the same size budget

So a `.docx` citation that says "page 7" means "the 7th section as this parser
divided it", NOT "what Word shows as page 7" — Word paginates at render time
from fonts, margins and printer metrics that are not in the file. The division
IS deterministic and reproducible, which is what a citation needs to be; it just
is not Word's. Properly fixing the label means giving `Citation` a format-aware
locator instead of a bare `page`, which is a schema change and belongs with the
span-citation work, not here.
"""
import os
import re
from collections.abc import Callable
from functools import lru_cache
from pathlib import Path

import fitz  # PyMuPDF

from backend.config import settings
from backend.models import UNIT_PAGE, UNIT_SECTION, UNIT_SHEET, UNIT_SLIDE

# Size budget for formats with no page concept of their own. Roughly a dense
# page of prose, and comfortably larger than one chunk, so a section still
# splits into several chunks rather than the two boundaries fighting.
SECTION_CHARS = 3000


class UnsupportedDocument(ValueError):
    """Raised for a file extension no extractor handles."""


# ——— PDF ———


def _pdf_pages(path: Path) -> list[str]:
    """Extract a PDF page by page, falling back to OCR on imaged pages.

    MEASURED AND REJECTED: `page.get_text(sort=True)`. It was expected to fix
    multi-column reading order; it does the opposite. `sort` orders text blocks
    by vertical position first, so on a two-column page it reads *across* both
    columns line by line and interleaves them:

        default      "...vacancy fell to 3.1% in the third quarter..."
        sort=True    "...vacancy fell to    Rail freight volumes tell a
                      3.1% in the third     different story from the..."

    PyMuPDF's default order follows the content stream, which for a well-formed
    PDF already runs column by column. Verified on a two-column brief and on the
    side-by-side result tables of a real LaTeX paper; on the 30-question eval it
    moved 3 questions (2 worse, 1 better) — noise — while visibly destroying
    two-column text. It would only help on a PDF whose content stream is already
    scrambled, and that is not the common case.
    """
    doc = fitz.open(str(path))
    try:
        return [_pdf_page_text(page) for page in doc]
    finally:
        doc.close()


def _pdf_page_text(page: fitz.Page) -> str:
    text = page.get_text()
    if text.strip():
        return text
    # No text layer on this page: a scan, or a page that is purely an image.
    # OCR only these pages — running it on a page that already has text would be
    # slower and worse than the text the PDF already carries.
    return _ocr_page(page)


# Standard install locations, checked when nothing is configured. PyMuPDF shells
# out to Tesseract and needs to be told where the language data is; it does not
# search for it. Ordered most- to least-likely.
_TESSDATA_CANDIDATES = (
    r"C:\Program Files\Tesseract-OCR\tessdata",
    r"C:\Program Files (x86)\Tesseract-OCR\tessdata",
    "/usr/share/tesseract-ocr/5/tessdata",
    "/usr/share/tesseract-ocr/4.00/tessdata",
    "/usr/share/tessdata",
    "/opt/homebrew/share/tessdata",
    "/usr/local/share/tessdata",
)


@lru_cache(maxsize=1)
def _tessdata() -> str | None:
    """Locate Tesseract's language data, or None if OCR is unavailable here.

    Explicit config wins, then the conventional environment variable, then the
    standard install paths — so it works out of the box on a normal install but
    can still be pinned for a container that puts tessdata elsewhere.
    """
    for candidate in (settings.tessdata_prefix, os.environ.get("TESSDATA_PREFIX", ""),
                      *_TESSDATA_CANDIDATES):
        if candidate and Path(candidate).is_dir():
            return candidate
    return None


@lru_cache(maxsize=1)
def _tesseract_works() -> bool:
    """Probe once whether Tesseract can actually run here.

    Cached because it shells out to the binary, and whether a binary exists does
    not change while the process runs — probing per scanned page would pay that
    cost on every page of every scan.

    Note what this does NOT test: passing `tessdata=None` is not the same as
    "no OCR". Tesseract falls back to its own compiled-in data location, which
    is exactly how a Linux package install works — so a machine with no
    configured tessdata path can still OCR perfectly well. Disabling OCR is
    `settings.ocr_enabled`, below, not the absence of a path.
    """
    doc = fitz.open()
    try:
        doc.new_page().get_textpage_ocr(tessdata=_tessdata())
        return True
    except (RuntimeError, TypeError):
        return False
    finally:
        doc.close()


def ocr_available() -> bool:
    """Whether a scanned page will actually be OCR'd: switched on AND working."""
    return settings.ocr_enabled and _tesseract_works()


def _ocr_page(page: fitz.Page) -> str:
    """OCR one page, or return "" if OCR is unavailable.

    Returning "" rather than raising is deliberate: unavailable OCR must degrade
    a scanned page to "no extractable text" — the same 422 the upload already
    produces — instead of turning every mixed document into a hard failure
    because one page happened to be an image.
    """
    if not ocr_available():
        return ""
    try:
        return page.get_text(textpage=page.get_textpage_ocr(full=False, tessdata=_tessdata()))
    except (RuntimeError, TypeError):
        return ""


# ——— Word ———


def _docx_pages(path: Path) -> list[str]:
    """Extract a .docx as sections, split on explicit page breaks where present.

    Walks the body in document order so a table stays between the paragraphs
    that introduce and follow it — python-docx's `.paragraphs` and `.tables`
    would each give a flat list with the interleaving lost.
    """
    import docx
    from docx.document import Document as DocxDocument
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = docx.Document(str(path))
    parts: list[str] = []
    for block in _iter_docx_blocks(doc, DocxDocument, Paragraph, Table):
        if isinstance(block, Paragraph):
            if _paragraph_has_page_break(block):
                parts.append(_PAGE_BREAK)
            if block.text.strip():
                parts.append(block.text)
        else:
            parts.append(_table_to_markdown([[c.text for c in row.cells] for row in block.rows]))

    return _split_sections(parts)


_PAGE_BREAK = "\x00PAGEBREAK\x00"  # sentinel; never appears in real document text


def _iter_docx_blocks(doc, DocxDocument, Paragraph, Table):
    """Yield paragraphs and tables in the order they appear in the document."""
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag.endswith("}p"):
            yield Paragraph(child, doc)
        elif child.tag.endswith("}tbl"):
            yield Table(child, doc)


def _paragraph_has_page_break(paragraph) -> bool:
    """True if the author put an explicit page break at the start of this run."""
    return any(
        br.get(f"{{{br.nsmap['w']}}}type") == "page"
        for run in paragraph.runs
        for br in run._element.findall(f"{{{run._element.nsmap['w']}}}br")
    )


# ——— PowerPoint ———


def _pptx_pages(path: Path) -> list[str]:
    """One entry per slide: title, body text, tables, and the speaker notes.

    Notes are included because that is often where the actual argument lives —
    the slide says "Q3 Revenue" and the notes say why it moved.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                parts.append(
                    _table_to_markdown([[c.text for c in row.cells] for row in shape.table.rows])
                )
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(f"Speaker notes: {notes}")
        slides.append("\n\n".join(parts))
    return slides


# ——— Excel ———


def _xlsx_pages(path: Path) -> list[str]:
    """One entry per worksheet, rendered as a Markdown table.

    `data_only=True` reads the cached values Excel stored for formula cells
    rather than the formula source — "=B2*C2" is not an answer to a question
    about a number. Openpyxl cannot evaluate formulas, so a file written by a
    tool that never cached values yields None there; that is a real limitation
    of the format, not something to paper over.
    """
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    try:
        sheets: list[str] = []
        for ws in wb.worksheets:
            rows = [
                ["" if cell is None else str(cell) for cell in row]
                for row in ws.iter_rows(values_only=True)
            ]
            rows = [r for r in rows if any(c.strip() for c in r)]
            body = _table_to_markdown(rows) if rows else ""
            sheets.append(f"Sheet: {ws.title}\n\n{body}".strip())
        return sheets
    finally:
        wb.close()


# ——— Plain text ———


def _text_pages(path: Path) -> list[str]:
    """Read as UTF-8, tolerating the odd bad byte rather than failing the upload."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return _split_sections(text.split("\n\n"))


# ——— shared helpers ———


def _table_to_markdown(rows: list[list[str]]) -> str:
    """Render a table so each number keeps its row and column label.

    This is the whole point of handling tables separately. Flattened, a fee
    table reads as "Standard 25 180,000 Professional 100 480,000" — the numbers
    are present but unattached to anything, and no amount of retrieval quality
    makes that answerable. As a Markdown table each cell still sits under its
    header, which is a form both the embedder and the model can use.
    """
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    padded = [[_cell(c) for c in r] + [""] * (width - len(r)) for r in rows]
    header, *body = padded
    lines = ["| " + " | ".join(header) + " |", "|" + "---|" * width]
    lines += ["| " + " | ".join(row) + " |" for row in body]
    return "\n".join(lines)


def _cell(text: str) -> str:
    """Flatten a cell to one line and escape the pipe that would break the table."""
    return re.sub(r"\s+", " ", str(text)).replace("|", r"\|").strip()


def _split_sections(parts: list[str]) -> list[str]:
    """Group parts into sections of at most SECTION_CHARS, honouring page breaks.

    A single part longer than the budget is never cut — chunking already handles
    oversized text, and slicing a table in half here would undo the work of
    building it.
    """
    sections: list[str] = []
    current: list[str] = []
    size = 0

    def flush() -> None:
        nonlocal current, size
        if current:
            sections.append("\n\n".join(current))
        current, size = [], 0

    for part in parts:
        if part == _PAGE_BREAK:
            flush()
            continue
        if not part.strip():
            continue
        if size and size + len(part) > SECTION_CHARS:
            flush()
        current.append(part)
        size += len(part)

    flush()
    return sections


# ——— dispatch ———

EXTRACTORS: dict[str, Callable[[Path], list[str]]] = {
    ".pdf": _pdf_pages,
    ".docx": _docx_pages,
    ".pptx": _pptx_pages,
    ".xlsx": _xlsx_pages,
    ".txt": _text_pages,
    ".md": _text_pages,
}

# What one entry in each format's list actually is, so a citation can say
# "slide 3" or "sheet 2" rather than calling everything a page. See the module
# docstring for why `.docx` is a section and not a page.
UNITS: dict[str, str] = {
    ".pdf": UNIT_PAGE,
    ".docx": UNIT_SECTION,
    ".pptx": UNIT_SLIDE,
    ".xlsx": UNIT_SHEET,
    ".txt": UNIT_SECTION,
    ".md": UNIT_SECTION,
}

SUPPORTED_EXTENSIONS = frozenset(EXTRACTORS)


def unit_for(doc_path: str | Path) -> str:
    """The name for one addressable location in this document's format."""
    return UNITS.get(Path(doc_path).suffix.lower(), UNIT_PAGE)


def extract_pages(doc_path: str | Path) -> list[str]:
    """Return the document's text, one entry per page-like unit (0-indexed).

    Raises UnsupportedDocument for an extension with no extractor. Anything the
    underlying library raises on a corrupt file propagates — the upload route
    turns that into a 422 and rolls the stored file back.
    """
    path = Path(doc_path)
    extractor = EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        raise UnsupportedDocument(
            f"unsupported file type '{path.suffix}' "
            f"(supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))})"
        )
    return extractor(path)
