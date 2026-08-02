"""Shared fixtures — and, first, the isolation that makes this suite safe to run.

The app resolves BOTH of its storage locations at import time:
  * `backend.db.engine` is built from `settings.database_url`, and
  * `backend.library.DATA_ROOT` is a module constant.

So the environment MUST be redirected before `backend.*` is imported for the
first time, which is why the os.environ writes below sit at module top level and
run before the imports rather than inside a fixture. A hard assertion then
proves both point somewhere under a throwaway temp directory — if a future
change breaks the redirection, the suite refuses to run instead of quietly
writing into the developer's real library.
"""
import os
import shutil
import tempfile
from pathlib import Path

import pytest

# ——— redirect storage BEFORE importing anything from backend ———

_TMP = Path(tempfile.mkdtemp(prefix="scholar-tests-"))
os.environ["DATABASE_URL"] = f"sqlite:///{(_TMP / 'test.db').as_posix()}"
os.environ["SECRET_KEY"] = "test-only-secret-not-used-anywhere-real"
os.environ["ACCESS_TOKEN_EXPIRE_MINUTES"] = "30"

import fitz  # noqa: E402  (PDF fixtures are built, not committed as binaries)
from fastapi.testclient import TestClient  # noqa: E402

from backend import library  # noqa: E402
from backend.api import app  # noqa: E402
from backend.config import settings  # noqa: E402
from backend.db import engine, init_db  # noqa: E402

library.DATA_ROOT = _TMP / "users"

# ——— the guard: prove we are not pointed at real data ———

_REPO_DATA = (Path(__file__).resolve().parents[2] / "data").resolve()

assert _TMP.resolve() != _REPO_DATA, "temp dir collided with the repo data dir"
assert str(library.DATA_ROOT.resolve()).startswith(str(_TMP.resolve())), (
    f"library.DATA_ROOT escaped the sandbox: {library.DATA_ROOT}"
)
assert str(_TMP.resolve()) in settings.database_url.replace("/", os.sep), (
    f"DATABASE_URL escaped the sandbox: {settings.database_url}"
)
assert not str(engine.url).endswith("scholar.db"), (
    f"engine is bound to the real database: {engine.url}"
)


def pytest_sessionfinish(session, exitstatus):
    """Delete the throwaway tree once the whole run is over."""
    shutil.rmtree(_TMP, ignore_errors=True)


# ——— fixtures ———


@pytest.fixture(scope="session", autouse=True)
def _database():
    init_db()


@pytest.fixture
def client() -> TestClient:
    """A TestClient with a clean database and empty user storage.

    Built WITHOUT the `with` form on purpose: entering the context manager runs
    the app's lifespan, which spawns the model-warming thread and tries to reach
    Ollama. Tests call init_db themselves and never need that.
    """
    _reset_state()
    return TestClient(app)


def _reset_state() -> None:
    """Truncate every table and wipe per-user storage between tests.

    Also clears the two process-global caches. Without this the rate limiters
    accumulate across the whole session and later tests start getting 429s that
    have nothing to do with what they are testing — which makes results depend
    on test ORDER, the least debuggable kind of failure.
    """
    from sqlmodel import Session, SQLModel, delete

    from backend.ratelimit import ask_limiter, upload_limiter

    with Session(engine) as session:
        for table in reversed(SQLModel.metadata.sorted_tables):
            session.exec(delete(table))
        session.commit()
    shutil.rmtree(library.DATA_ROOT, ignore_errors=True)
    library._retrievers.clear()  # the per-user Retriever cache is process-global
    ask_limiter.reset()
    upload_limiter.reset()


@pytest.fixture
def alice(client: TestClient) -> TestClient:
    """A registered, logged-in client. Registration sets both auth cookies."""
    r = client.post(
        "/auth/register",
        json={"email": "alice@example.com", "password": "validpassword123"},
    )
    assert r.status_code == 201, r.text
    return client


def csrf(client: TestClient) -> dict[str, str]:
    """The double-submit header for an authenticated client's unsafe request."""
    return {"X-CSRF-Token": client.cookies.get("csrf_token", "")}


@pytest.fixture
def other_client() -> TestClient:
    """A SECOND client against the same app — a different browser, in effect.

    Shares the app and database with `client` but has its own cookie jar, which
    is what makes cross-tenant isolation testable.
    """
    return TestClient(app)


# ——— PDF fixtures, generated rather than committed ———


def make_pdf(path: Path, pages: list[str]) -> Path:
    """Write a real, well-formed PDF with the given text on each page."""
    path.parent.mkdir(parents=True, exist_ok=True)
    doc = fitz.open()
    for text in pages:
        page = doc.new_page()
        page.insert_textbox(fitz.Rect(56, 56, 556, 736), text, fontsize=11)
    doc.save(str(path))
    doc.close()
    return path


BODY = (
    "Retrieval-augmented generation combines a parametric model with a "
    "non-parametric memory. The retriever selects passages and the "
    "generator conditions on them. "
) * 12


@pytest.fixture
def text_pdf(tmp_path: Path) -> bytes:
    """A small but genuinely indexable PDF (enough text to survive chunking)."""
    return make_pdf(tmp_path / "small.pdf", [BODY, BODY]).read_bytes()


@pytest.fixture
def scanned_pdf(tmp_path: Path) -> bytes:
    """A genuinely scanned PDF: rendered, rasterised, text layer discarded.

    Built rather than checked in so the expected text is known ground truth to
    compare OCR output against.
    """
    body = "Force Majeure clause seven point two. " + BODY[:400]
    src = fitz.open()
    page = src.new_page()
    page.insert_textbox(fitz.Rect(50, 50, 550, 740), body, fontsize=14)
    pixmap = page.get_pixmap(dpi=200)

    scan = fitz.open()
    scanned = scan.new_page(width=page.rect.width, height=page.rect.height)
    scanned.insert_image(scanned.rect, pixmap=pixmap)
    path = tmp_path / "scanned.pdf"
    scan.save(str(path))
    scan.close()
    src.close()
    return path.read_bytes()


@pytest.fixture
def docx_bytes(tmp_path: Path) -> bytes:
    import docx

    doc = docx.Document()
    doc.add_paragraph("Master Services Agreement")
    doc.add_paragraph(BODY)
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Tier"
    table.rows[0].cells[1].text = "Annual Fee"
    table.rows[1].cells[0].text = "Enterprise"
    table.rows[1].cells[1].text = "960,000"
    path = tmp_path / "agreement.docx"
    doc.save(str(path))
    return path.read_bytes()


@pytest.fixture
def xlsx_bytes(tmp_path: Path) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Fees"
    ws.append(["Tier", "Included Seats", "Annual Fee"])
    ws.append(["Standard", 25, 180000])
    ws.append(["Enterprise", "unlimited", 960000])
    path = tmp_path / "fees.xlsx"
    wb.save(str(path))
    return path.read_bytes()


@pytest.fixture
def pptx_bytes(tmp_path: Path) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Review"
    slide.notes_slide.notes_text_frame.text = BODY
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path.read_bytes()
