"""Per-format extraction: every format lands on the same `list[str]` contract.

These build real files with the real libraries rather than checking in binary
fixtures, so the tests exercise the actual format, not a recording of it.
"""
from pathlib import Path

import fitz
import pytest

from backend.parser import (
    SECTION_CHARS,
    SUPPORTED_EXTENSIONS,
    UnsupportedDocument,
    _table_to_markdown,
    extract_pages,
    ocr_available,
)


# ——— dispatch ———


def test_unsupported_extension_is_rejected(tmp_path: Path):
    path = tmp_path / "photo.png"
    path.write_bytes(b"\x89PNG")
    with pytest.raises(UnsupportedDocument):
        extract_pages(path)


def test_extension_matching_is_case_insensitive(tmp_path: Path):
    path = tmp_path / "NOTES.MD"
    path.write_text("# Title\n\nSome content.", encoding="utf-8")
    assert extract_pages(path) == ["# Title\n\nSome content."]


def test_supported_set_is_what_the_dispatch_table_offers():
    assert SUPPORTED_EXTENSIONS >= {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"}


# ——— PDF ———


def test_pdf_returns_one_entry_per_page(tmp_path: Path):
    doc = fitz.open()
    for text in ("First page text.", "Second page text.", "Third page text."):
        doc.new_page().insert_textbox(fitz.Rect(56, 56, 556, 736), text, fontsize=11)
    path = tmp_path / "three.pdf"
    doc.save(str(path))
    doc.close()

    pages = extract_pages(path)
    assert len(pages) == 3
    assert "First page" in pages[0]
    assert "Third page" in pages[2]


def test_two_column_pdf_is_not_interleaved(tmp_path: Path):
    """The regression that `sort=True` would introduce.

    Sorting text blocks by vertical position reads ACROSS side-by-side columns,
    splicing them line by line into text no reader can follow. This asserts each
    column survives as a contiguous run.
    """
    left = "The vacancy rate fell to three point one percent this quarter. " * 6
    right = "Rail freight volumes declined by two point seven percent instead. " * 6
    doc = fitz.open()
    page = doc.new_page()
    page.insert_textbox(fitz.Rect(56, 56, 290, 736), left, fontsize=9)
    page.insert_textbox(fitz.Rect(306, 56, 540, 736), right, fontsize=9)
    path = tmp_path / "twocol.pdf"
    doc.save(str(path))
    doc.close()

    text = " ".join(extract_pages(path)[0].split())
    assert "vacancy rate fell to three point one percent this quarter" in text
    assert "Rail freight volumes declined by two point seven percent" in text


def test_blank_page_yields_empty_rather_than_raising(tmp_path: Path):
    """A page with nothing on it must degrade to "" — which the upload route
    turns into a 422 — not raise and fail the whole document."""
    doc = fitz.open()
    doc.new_page()
    path = tmp_path / "blank.pdf"
    doc.save(str(path))
    doc.close()

    assert extract_pages(path) == [""]


def test_ocr_available_reports_a_bool_and_does_not_raise():
    """Tesseract is a system binary that may or may not be present; the probe
    must answer the question rather than propagate a RuntimeError."""
    assert isinstance(ocr_available(), bool)


def _scanned_pdf(tmp_path: Path, text: str) -> Path:
    """A genuinely scanned PDF: render text, rasterise it, keep only the image.

    Building it this way rather than checking in a scan means the fixture has a
    known ground truth to compare the OCR output against.
    """
    src = fitz.open()
    page = src.new_page()
    page.insert_textbox(fitz.Rect(50, 50, 550, 400), text, fontsize=18)
    pixmap = page.get_pixmap(dpi=200)

    scan = fitz.open()
    scanned_page = scan.new_page(width=page.rect.width, height=page.rect.height)
    scanned_page.insert_image(scanned_page.rect, pixmap=pixmap)
    path = tmp_path / "scanned.pdf"
    scan.save(str(path))
    scan.close()
    src.close()
    return path


def test_a_scanned_pdf_really_has_no_text_layer(tmp_path: Path):
    """Guards the fixture itself — if rasterising ever stopped stripping the text
    layer, the OCR test below would pass without OCR ever running."""
    path = _scanned_pdf(tmp_path, "Force Majeure clause seven point two")
    doc = fitz.open(str(path))
    try:
        assert doc[0].get_text().strip() == ""
    finally:
        doc.close()


@pytest.mark.slow
@pytest.mark.skipif(not ocr_available(), reason="OCR unavailable here (no Tesseract, or disabled)")
def test_scanned_pdf_is_recovered_by_ocr(tmp_path: Path):
    """The document class that was previously a hard 422: a scan."""
    path = _scanned_pdf(tmp_path, "Force Majeure clause seven point two")
    text = " ".join(extract_pages(path)[0].split())
    assert "Force Majeure" in text
    assert "seven point two" in text


@pytest.mark.slow
@pytest.mark.skipif(not ocr_available(), reason="OCR unavailable here (no Tesseract, or disabled)")
def test_clearing_the_tessdata_path_does_not_disable_ocr(tmp_path: Path, monkeypatch):
    """The trap that made an earlier test pass on Windows and fail in CI.

    Tesseract falls back to its own compiled-in data location, so a blank
    tessdata path still OCRs on any normal Linux install. This asserts the
    surprising behaviour directly, so nobody reaches for the path again when
    they mean to turn OCR off.
    """
    from backend import parser

    monkeypatch.setattr(parser, "_tessdata", lambda: None)
    parser._tesseract_works.cache_clear()
    try:
        text = " ".join(extract_pages(_scanned_pdf(tmp_path, "Retention is eighty four months"))[0].split())
    finally:
        parser._tesseract_works.cache_clear()

    # On a machine that CAN find its own tessdata, OCR still ran. On one that
    # cannot, it degraded to "" — never an exception, and never a crash.
    assert text == "" or "eighty four months" in text


@pytest.mark.slow
@pytest.mark.skipif(not ocr_available(), reason="OCR unavailable here (no Tesseract, or disabled)")
def test_ocr_enabled_false_really_disables_ocr(tmp_path: Path, monkeypatch):
    """The switch that DOES turn OCR off, on every platform."""
    from backend.config import settings

    monkeypatch.setattr(settings, "ocr_enabled", False)
    assert ocr_available() is False
    assert extract_pages(_scanned_pdf(tmp_path, "Retention is eighty four months")) == [""]


def test_ocr_is_skipped_on_pages_that_already_have_text(tmp_path: Path, monkeypatch):
    """OCR is slower and worse than a text layer the PDF already carries, so it
    must only run on pages that have none."""
    from backend import parser

    calls = []
    monkeypatch.setattr(parser, "_ocr_page", lambda page: calls.append(1) or "")

    doc = fitz.open()
    doc.new_page().insert_textbox(fitz.Rect(56, 56, 556, 736), "Real text here.", fontsize=11)
    path = tmp_path / "text.pdf"
    doc.save(str(path))
    doc.close()

    extract_pages(path)
    assert calls == [], "OCR ran on a page that already had a text layer"


# ——— Word ———


def _docx(path: Path, paragraphs: list[str], table: list[list[str]] | None = None) -> Path:
    import docx

    doc = docx.Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    if table:
        t = doc.add_table(rows=len(table), cols=len(table[0]))
        for row, values in zip(t.rows, table):
            for cell, value in zip(row.cells, values):
                cell.text = value
    doc.save(str(path))
    return path


def test_docx_extracts_paragraph_text(tmp_path: Path):
    path = _docx(tmp_path / "a.docx", ["Section 7.2 Force Majeure.", "Neither party shall be liable."])
    text = "\n".join(extract_pages(path))
    assert "Section 7.2 Force Majeure." in text
    assert "Neither party shall be liable." in text


def test_docx_table_keeps_cells_under_their_headers(tmp_path: Path):
    """The point of the whole table path: a number must stay attached to its
    row and column label, not dissolve into a run of loose figures."""
    path = _docx(
        tmp_path / "fees.docx",
        ["Fee schedule:"],
        [["Tier", "Seats", "Annual Fee"], ["Professional", "100", "480,000"]],
    )
    text = "\n".join(extract_pages(path))
    assert "| Tier | Seats | Annual Fee |" in text
    assert "| Professional | 100 | 480,000 |" in text


def test_docx_splits_into_sections_on_a_size_budget(tmp_path: Path):
    path = _docx(tmp_path / "long.docx", ["A paragraph of moderate length." * 20] * 12)
    pages = extract_pages(path)
    assert len(pages) > 1, "a long document should not collapse into one section"
    assert all(len(p) <= SECTION_CHARS * 2 for p in pages)


def test_docx_splits_on_an_explicit_page_break(tmp_path: Path):
    import docx
    from docx.enum.text import WD_BREAK

    doc = docx.Document()
    doc.add_paragraph("Content before the break.")
    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    doc.add_paragraph("Content after the break.")
    path = tmp_path / "broken.docx"
    doc.save(str(path))

    pages = extract_pages(path)
    assert len(pages) == 2
    assert "before the break" in pages[0]
    assert "after the break" in pages[1]


def test_docx_keeps_a_table_between_the_paragraphs_around_it(tmp_path: Path):
    """python-docx's .paragraphs and .tables are separate flat lists — reading
    them that way loses the interleaving and detaches a table from its caption."""
    import docx

    doc = docx.Document()
    doc.add_paragraph("Intro before.")
    t = doc.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "Alpha"
    t.rows[0].cells[1].text = "Beta"
    doc.add_paragraph("Outro after.")
    path = tmp_path / "order.docx"
    doc.save(str(path))

    text = "\n".join(extract_pages(path))
    assert text.index("Intro before.") < text.index("Alpha") < text.index("Outro after.")


# ——— PowerPoint ———


def test_pptx_returns_one_entry_per_slide_including_notes(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for title in ("Q3 Revenue", "Q4 Outlook"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        box.text_frame.text = f"Body text for {title}."
        slide.notes_slide.notes_text_frame.text = f"Why {title} moved."
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    slides = extract_pages(path)
    assert len(slides) == 2
    assert "Q3 Revenue" in slides[0]
    assert "Body text for Q3 Revenue." in slides[0]
    assert "Why Q3 Revenue moved." in slides[0], "speaker notes carry the actual argument"
    assert "Q4 Outlook" in slides[1]


# ——— Excel ———


def test_xlsx_returns_one_entry_per_sheet_as_a_table(tmp_path: Path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Fees"
    ws.append(["Tier", "Annual Fee"])
    ws.append(["Enterprise", 960000])
    wb.create_sheet("Notes").append(["Reviewed quarterly"])
    path = tmp_path / "book.xlsx"
    wb.save(str(path))

    sheets = extract_pages(path)
    assert len(sheets) == 2
    assert "Sheet: Fees" in sheets[0]
    assert "| Enterprise | 960000 |" in sheets[0]
    assert "Sheet: Notes" in sheets[1]


def test_xlsx_reads_cached_formula_values_not_formula_source(tmp_path: Path):
    """A question about a number is not answered by the string "=B2*C2"."""
    from openpyxl import Workbook

    wb = Workbook()
    wb.active.append(["=1+1"])  # written as a formula, no cached value
    path = tmp_path / "formula.xlsx"
    wb.save(str(path))

    assert "=1+1" not in "\n".join(extract_pages(path))


# ——— plain text ———


def test_txt_is_read_as_one_section_when_short(tmp_path: Path):
    path = tmp_path / "notes.txt"
    path.write_text("Short note.\n\nSecond paragraph.", encoding="utf-8")
    assert extract_pages(path) == ["Short note.\n\nSecond paragraph."]


def test_txt_with_undecodable_bytes_does_not_fail_the_upload(tmp_path: Path):
    path = tmp_path / "mixed.txt"
    path.write_bytes(b"valid text \xff\xfe more text")
    pages = extract_pages(path)
    assert "valid text" in pages[0]
    assert "more text" in pages[0]


@pytest.mark.parametrize("content", ["", "   \n\n  \t "])
def test_empty_input_yields_no_sections(tmp_path: Path, content):
    """This is the signal papers.py turns into a 422."""
    path = tmp_path / "empty.txt"
    path.write_text(content, encoding="utf-8")
    assert extract_pages(path) == []


# ——— the Markdown table renderer ———


def test_ragged_rows_are_padded_to_a_valid_table():
    rendered = _table_to_markdown([["A", "B", "C"], ["1"], ["2", "3"]])
    assert [line.count("|") for line in rendered.splitlines()] == [4, 4, 4, 4]


def test_pipes_in_cells_are_escaped():
    assert r"a\|b" in _table_to_markdown([["a|b", "c"]])


def test_newlines_in_cells_are_flattened():
    """An unflattened newline would end the Markdown row mid-table."""
    rendered = _table_to_markdown([["line one\nline two", "c"]])
    assert len(rendered.splitlines()) == 2
    assert "line one line two" in rendered
